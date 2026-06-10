from pptx import Presentation

def ppt_export_agent(presentation_text: str, filename: str = "report.pptx"):

    prs = Presentation()

    slides = presentation_text.split("# Slide")

    for slide in slides:
        if slide.strip() == "":
            continue

        slide_layout = prs.slides.add_slide(prs.slide_layouts[1])

        parts = slide.split("\n")
        title = parts[0].strip()
        content = "\n".join(parts[1:])

        slide_layout.shapes.title.text = title
        slide_layout.placeholders[1].text = content

    prs.save(filename)

    return filename